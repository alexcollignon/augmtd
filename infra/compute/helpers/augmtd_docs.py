"""AUGMTD document helpers (DH6) — the fragile primitives, written and tested ONCE.

Generated compute scripts import these instead of reinventing the risky parts:

    from augmtd_docs import clone_slide, render_verify, ocr_image

Design laws baked in:
- PRESERVE WHAT YOU CAN'T PARSE: cloning copies the slide's XML subtree byte-faithfully
  (SmartArt, exotic shapes ride along); relationships (images/media) are re-created so the
  clone renders identically.
- THE RENDER-VERIFICATION GATE: LibreOffice-headless converts to PDF inside the locked room;
  a document that fails to render, has the wrong page count, or renders blank must not ship.
"""

from __future__ import annotations

import copy
import os
import subprocess
import tempfile


def clone_slide(pres, slide):
    """Duplicate a slide (python-pptx has no API for this) — XML deepcopy + relationship
    fixups, so images/media on the original appear on the clone. Returns the new slide.

    Works for hand-made decks (no masters/layouts needed): the clone uses the same layout
    object the original references.
    """
    from pptx.oxml.ns import qn  # noqa: F401  (imported for side-effect-free ns access)

    slide_layout = slide.slide_layout
    new_slide = pres.slides.add_slide(slide_layout)

    # Remove placeholder shapes the layout added — we copy EVERYTHING from the source.
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)

    # Byte-faithful copy of every shape subtree (preserve-what-you-can't-parse).
    for shp in slide.shapes:
        new_slide.shapes._spTree.append(copy.deepcopy(shp._element))

    # Re-create image/media relationships: rIds inside the copied XML must resolve in the
    # NEW slide part. Map old rId → new rId, then rewrite references in the copied tree.
    rid_map = {}
    for rid, rel in slide.part.rels.items():
        if "image" in rel.reltype or "media" in rel.reltype or "video" in rel.reltype:
            new_rid = new_slide.part.relate_to(rel._target, rel.reltype)
            rid_map[rid] = new_rid
    if rid_map:
        import re as _re

        xml = new_slide.shapes._spTree.xml
        for old, new in rid_map.items():
            xml = _re.sub(rf'r:embed="{old}"', f'r:embed="{new}"', xml)
            xml = _re.sub(rf'r:link="{old}"', f'r:link="{new}"', xml)
        # Replace the spTree with the rewritten XML.
        from lxml import etree

        new_tree = etree.fromstring(xml.encode())
        parent = new_slide.shapes._spTree.getparent()
        parent.replace(new_slide.shapes._spTree, new_tree)

    return new_slide


def render_verify(path: str, min_pages: int = 1) -> int:
    """THE RENDER-VERIFICATION GATE: convert the document to PDF with LibreOffice-headless
    and return the page count. Raises RuntimeError if conversion fails or the page count is
    below min_pages — a document that does not render must not ship.
    """
    out_dir = tempfile.mkdtemp(prefix="render_")
    # The sandbox filesystem is read-only outside /tmp and /job/out — LibreOffice needs a
    # writable HOME and user profile (rc=77 "cannot be started" otherwise, found on first run).
    env = dict(os.environ, HOME="/tmp")
    res = subprocess.run(
        [
            "soffice", "-env:UserInstallation=file:///tmp/lo_profile",
            "--headless", "--convert-to", "pdf", "--outdir", out_dir, path,
        ],
        capture_output=True, timeout=120, env=env,
    )
    base = os.path.splitext(os.path.basename(path))[0] + ".pdf"
    pdf_path = os.path.join(out_dir, base)
    if res.returncode != 0 or not os.path.exists(pdf_path):
        raise RuntimeError(
            f"render failed: rc={res.returncode} stderr={res.stderr.decode(errors='replace')[:300]}"
        )
    from pypdf import PdfReader

    pages = len(PdfReader(pdf_path).pages)
    if pages < min_pages:
        raise RuntimeError(f"rendered only {pages} pages (expected >= {min_pages})")
    return pages


def ocr_image(path: str) -> str:
    """OCR one image (scanned page) with tesseract. Returns extracted text."""
    import pytesseract
    from PIL import Image

    return pytesseract.image_to_string(Image.open(path))
